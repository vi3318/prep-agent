from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.chart import XL_CHART_TYPE
import os
import re
import requests
import threading

# --- Theme Constants ---
HEADER_COLOR = RGBColor(0, 70, 140)
BODY_COLOR = RGBColor(20, 20, 20)
FOOTER_COLOR = RGBColor(100, 100, 100)
LIGHT_BLUE = RGBColor(222, 240, 255)
DARKER_BLUE = RGBColor(180, 210, 240)
FOOTER_FONT_SIZE = Pt(14)
TITLE_FONT_SIZE = Pt(40)
SECTION_FONT_SIZE = Pt(28)
BODY_FONT_SIZE = Pt(18)
FONT_NAME = "Calibri"

SECTION_TITLES = [
    "Executive Summary",
    "Key Offerings & Business Segments",
    "Strategic Direction & Initiatives",
    "Market Positioning & Target Audience",
    "Additional Notes"
]

# --- Add modern Salesforce-style background ---
def add_modern_background(slide, slide_width, slide_height):
    # Solid light blue background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_width, slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_BLUE
    bg.line.fill.background()
    # Large semi-transparent oval/arc at the bottom
    arc = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        left=-int(slide_width * 0.3),
        top=int(slide_height * 0.65),
        width=int(slide_width * 1.6),
        height=int(slide_height * 0.8)
    )
    arc.fill.solid()
    arc.fill.fore_color.rgb = DARKER_BLUE
    arc.fill.transparency = 0.15
    arc.line.fill.background()

# --- Add footer ---
def add_footer(slide, slide_width, slide_height):
    footer = slide.shapes.add_textbox(0, slide_height - Inches(0.6), slide_width, Inches(0.3))
    p = footer.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Prepared by Client Prep Agent"
    run.font.size = FOOTER_FONT_SIZE
    run.font.color.rgb = FOOTER_COLOR
    run.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER

# --- Add Salesforce logo on title slide ---
def add_logo(slide, slide_width):
    logo_path = "logo.png"
    if os.path.exists(logo_path):
        slide.shapes.add_picture(
            logo_path,
            slide_width - Inches(1.5),
            Inches(0.3),
            Inches(1.2),
            Inches(1.2)
        )
    else:
        print("⚠️ Salesforce-logo.jpg not found")

# --- Robust summary parser ---
def parse_summary_sections(summary_text):
    # Accepts headings with or without bold, numbers, or colons
    # e.g., Executive Summary, 1. Executive Summary, **Executive Summary:**
    pattern = r"(?:\*\*)?(?:\d+\.\s*)?([A-Za-z &]+?)(?:\*\*)?:?\n"
    matches = list(re.finditer(pattern, summary_text))
    sections = {}
    if matches:
        for i, match in enumerate(matches):
            title = match.group(1).strip()
            start = match.end()
            end = matches[i+1].start() if i+1 < len(matches) else len(summary_text)
            body = summary_text[start:end].strip()
            sections[title] = body
    else:
        # Fallback: all content to first section
        sections[SECTION_TITLES[0]] = summary_text.strip()
    return sections

# --- Add title slide ---
def add_title_slide(prs, company_name, slide_width, slide_height):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_modern_background(slide, slide_width, slide_height)
    add_footer(slide, slide_width, slide_height)
    add_logo(slide, slide_width)
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.4),
        slide_width - Inches(2), Inches(2)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"{company_name} Summary Report"
    run.font.size = TITLE_FONT_SIZE
    run.font.bold = True
    run.font.color.rgb = HEADER_COLOR
    run.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER

# --- Add content slide ---
def add_content_slide(prs, title, body_text, slide_width, slide_height):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_modern_background(slide, slide_width, slide_height)
    add_footer(slide, slide_width, slide_height)
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1), slide_width - Inches(2), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = SECTION_FONT_SIZE
    run.font.bold = True
    run.font.color.rgb = HEADER_COLOR
    run.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER
    # Body
    body_box = slide.shapes.add_textbox(Inches(1), Inches(2.0), slide_width - Inches(2), Inches(4))
    tf = body_box.text_frame
    tf.word_wrap = True
    if body_text:
        for line in body_text.strip().split("\n"):
            clean = line.strip("-–•●* ").strip()
            if clean:
                p = tf.add_paragraph()
                p.text = f"• {clean}"
                p.level = 0
                p.font.size = BODY_FONT_SIZE
                p.font.color.rgb = BODY_COLOR
                p.font.name = FONT_NAME
    else:
        p = tf.add_paragraph()
        p.text = "(No content provided)"
        p.font.size = BODY_FONT_SIZE
        p.font.color.rgb = BODY_COLOR
        p.font.name = FONT_NAME

# --- Export summary to PPT (fixed 6 slides, robust parsing) ---
def export_summary_to_ppt(summary_text, filename, company_name="Company", company_url=None):
    prs = Presentation()
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    add_title_slide(prs, company_name, slide_width, slide_height)
    sections = parse_summary_sections(summary_text)
    for i, section_title in enumerate(SECTION_TITLES):
        body = ""
        # Try to match by exact or partial (case-insensitive)
        for k, v in sections.items():
            if section_title.lower() in k.lower():
                body = v
                break
        add_content_slide(prs, section_title, body, slide_width, slide_height)
    prs.save(filename)
    print(f"✅ Saved presentation: {filename}")

# --- Generate a .potx template with the same design ---
def generate_custom_template_potx(filename="Custom_Template.potx"):
    prs = Presentation()
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    add_title_slide(prs, "Sample Company", slide_width, slide_height)
    for i in range(5):
        add_content_slide(prs, SECTION_TITLES[i], "Section content here...", slide_width, slide_height)
    prs.save(filename)
    print(f"✅ Saved template: {filename}")

def add_financials_slide(prs, financials, slide_width, slide_height, company_name="Company"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_modern_background(slide, slide_width, slide_height)
    add_footer(slide, slide_width, slide_height)
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1), slide_width - Inches(2), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"Key Financials: {company_name}"
    run.font.size = SECTION_FONT_SIZE
    run.font.bold = True
    run.font.color.rgb = HEADER_COLOR
    run.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER
    # Table or bullets
    if financials:
        rows = len(financials)
        table = slide.shapes.add_table(rows, 2, Inches(2), Inches(2.2), Inches(6), Inches(0.6 * rows)).table
        for i, (k, v) in enumerate(financials.items()):
            table.cell(i, 0).text = k
            table.cell(i, 1).text = v
            for j in range(2):
                for run in table.cell(i, j).text_frame.paragraphs[0].runs:
                    run.font.size = BODY_FONT_SIZE
                    run.font.name = FONT_NAME
    else:
        body_box = slide.shapes.add_textbox(Inches(1), Inches(2.2), slide_width - Inches(2), Inches(2))
        tf = body_box.text_frame
        p = tf.paragraphs[0]
        p.text = "No financial data found."
        p.font.size = BODY_FONT_SIZE
        p.font.color.rgb = BODY_COLOR
        p.font.name = FONT_NAME

def add_swot_slide(prs, swot, slide_width, slide_height, company_name="Company"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_modern_background(slide, slide_width, slide_height)
    add_footer(slide, slide_width, slide_height)
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1), slide_width - Inches(2), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"SWOT Analysis: {company_name}"
    run.font.size = SECTION_FONT_SIZE
    run.font.bold = True
    run.font.color.rgb = HEADER_COLOR
    run.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER
    # SWOT sections
    y = 2.2
    for key in ["Strengths", "Weaknesses", "Opportunities", "Threats"]:
        box = slide.shapes.add_textbox(Inches(1), Inches(y), slide_width - Inches(2), Inches(0.7))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = key
        p.font.size = BODY_FONT_SIZE
        p.font.bold = True
        p.font.color.rgb = HEADER_COLOR
        p.font.name = FONT_NAME
        for item in swot.get(key, []):
            para = tf.add_paragraph()
            para.text = f"• {item}"
            para.font.size = BODY_FONT_SIZE
            para.font.color.rgb = BODY_COLOR
            para.font.name = FONT_NAME
        y += 0.7 + 0.2 * max(1, len(swot.get(key, [])))

def add_comparison_slide(prs, company1, company2, financials1, financials2, slide_width, slide_height):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_modern_background(slide, slide_width, slide_height)
    add_footer(slide, slide_width, slide_height)
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1), slide_width - Inches(2), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"Competitor Comparison"
    run.font.size = SECTION_FONT_SIZE
    run.font.bold = True
    run.font.color.rgb = HEADER_COLOR
    run.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER
    # Table
    keys = set(financials1.keys()) | set(financials2.keys())
    rows = len(keys) + 1
    table = slide.shapes.add_table(rows, 3, Inches(1.5), Inches(2.2), Inches(7), Inches(0.6 * rows)).table
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = company1
    table.cell(0, 2).text = company2
    for j in range(3):
        for run in table.cell(0, j).text_frame.paragraphs[0].runs:
            run.font.size = BODY_FONT_SIZE
            run.font.bold = True
            run.font.name = FONT_NAME
    for i, key in enumerate(keys, 1):
        table.cell(i, 0).text = key
        table.cell(i, 1).text = financials1.get(key, "N/A")
        table.cell(i, 2).text = financials2.get(key, "N/A")
        for j in range(3):
            for run in table.cell(i, j).text_frame.paragraphs[0].runs:
                run.font.size = BODY_FONT_SIZE
                run.font.name = FONT_NAME

def add_financials_bar_chart_slide(prs, financials, slide_width, slide_height, company_name="Company"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_modern_background(slide, slide_width, slide_height)
    add_footer(slide, slide_width, slide_height)
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1), slide_width - Inches(2), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"Key Financials (Bar Chart): {company_name}"
    run.font.size = SECTION_FONT_SIZE
    run.font.bold = True
    run.font.color.rgb = HEADER_COLOR
    run.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER
    # Prepare data
    chart_data = CategoryChartData()
    chart_data.categories = list(financials.keys())
    values = []
    for v in financials.values():
        try:
            # Remove currency symbols and text, keep numbers
            num = float(re.sub(r'[^\d.\-]', '', v.replace(',', '')))
            values.append(num)
        except Exception:
            values.append(0)
    chart_data.add_series('Value', values)
    # Add chart
    x, y, cx, cy = Inches(1.5), Inches(2.2), Inches(7), Inches(3)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart
    chart.has_legend = False
    chart.value_axis.has_major_gridlines = False
    chart.category_axis.tick_labels.font.size = Pt(14)
    chart.category_axis.tick_labels.font.name = FONT_NAME
    chart.value_axis.tick_labels.font.size = Pt(14)
    chart.value_axis.tick_labels.font.name = FONT_NAME
    chart.chart_title.text_frame.text = ""

def add_business_segments_pie_chart_slide(prs, segments, slide_width, slide_height, company_name="Company"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_modern_background(slide, slide_width, slide_height)
    add_footer(slide, slide_width, slide_height)
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1), slide_width - Inches(2), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"Business Segments (Pie Chart): {company_name}"
    run.font.size = SECTION_FONT_SIZE
    run.font.bold = True
    run.font.color.rgb = HEADER_COLOR
    run.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER
    # Prepare data
    chart_data = ChartData()
    chart_data.categories = list(segments.keys())
    chart_data.add_series('Segments', list(segments.values()))
    # Add chart
    x, y, cx, cy = Inches(2), Inches(2.2), Inches(6), Inches(3.5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
    ).chart
    chart.has_legend = True
    chart.legend.position = 2  # right
    chart.legend.font.size = Pt(14)
    chart.legend.font.name = FONT_NAME
    chart.chart_title.text_frame.text = ""

def add_trends_slide(prs, trends, slide_width, slide_height, company_name="Company"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_modern_background(slide, slide_width, slide_height)
    add_footer(slide, slide_width, slide_height)
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1), slide_width - Inches(2), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"Recent Trends: {company_name}"
    run.font.size = SECTION_FONT_SIZE
    run.font.bold = True
    run.font.color.rgb = HEADER_COLOR
    run.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER
    # Trends bullets
    body_box = slide.shapes.add_textbox(Inches(1), Inches(2.2), slide_width - Inches(2), Inches(3))
    tf = body_box.text_frame
    tf.word_wrap = True
    if trends:
        for trend in trends:
            p = tf.add_paragraph()
            p.text = f"• {trend}"
            p.font.size = BODY_FONT_SIZE
            p.font.color.rgb = BODY_COLOR
            p.font.name = FONT_NAME
    else:
        p = tf.add_paragraph()
        p.text = "No recent trends detected."
        p.font.size = BODY_FONT_SIZE
        p.font.color.rgb = BODY_COLOR
        p.font.name = FONT_NAME

def add_red_flags_opportunities_slide(prs, data, slide_width, slide_height, company_name="Company"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_modern_background(slide, slide_width, slide_height)
    add_footer(slide, slide_width, slide_height)
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1), slide_width - Inches(2), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"Red Flags & Opportunities: {company_name}"
    run.font.size = SECTION_FONT_SIZE
    run.font.bold = True
    run.font.color.rgb = HEADER_COLOR
    run.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER
    # Red Flags
    y = 2.2
    box = slide.shapes.add_textbox(Inches(1), Inches(y), slide_width - Inches(2), Inches(1))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "Red Flags"
    p.font.size = BODY_FONT_SIZE
    p.font.bold = True
    p.font.color.rgb = RGBColor(200, 40, 40)
    p.font.name = FONT_NAME
    for item in data.get('Red Flags', []):
        para = tf.add_paragraph()
        para.text = f"• {item}"
        para.font.size = BODY_FONT_SIZE
        para.font.color.rgb = BODY_COLOR
        para.font.name = FONT_NAME
    # Opportunities
    y += 1.0 + 0.2 * max(1, len(data.get('Red Flags', [])))
    box = slide.shapes.add_textbox(Inches(1), Inches(y), slide_width - Inches(2), Inches(1))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "Opportunities"
    p.font.size = BODY_FONT_SIZE
    p.font.bold = True
    p.font.color.rgb = RGBColor(40, 140, 40)
    p.font.name = FONT_NAME
    for item in data.get('Opportunities', []):
        para = tf.add_paragraph()
        para.text = f"• {item}"
        para.font.size = BODY_FONT_SIZE
        para.font.color.rgb = BODY_COLOR
        para.font.name = FONT_NAME

def add_timeline_slide(prs, events, slide_width, slide_height, company_name="Company"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_modern_background(slide, slide_width, slide_height)
    add_footer(slide, slide_width, slide_height)
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1), slide_width - Inches(2), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"Timeline: {company_name}"
    run.font.size = SECTION_FONT_SIZE
    run.font.bold = True
    run.font.color.rgb = HEADER_COLOR
    run.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER
    # Timeline events
    body_box = slide.shapes.add_textbox(Inches(1), Inches(2.2), slide_width - Inches(2), Inches(4))
    tf = body_box.text_frame
    tf.word_wrap = True
    if events:
        for year, desc in events:
            p = tf.add_paragraph()
            p.text = f"{year}: {desc}"
            p.font.size = BODY_FONT_SIZE
            p.font.color.rgb = BODY_COLOR
            p.font.name = FONT_NAME
    else:
        p = tf.add_paragraph()
        p.text = "No timeline events detected."
        p.font.size = BODY_FONT_SIZE
        p.font.color.rgb = BODY_COLOR
        p.font.name = FONT_NAME

# Example usage:
# add_financials_bar_chart_slide(prs, financials_dict, slide_width, slide_height, company_name)
# add_business_segments_pie_chart_slide(prs, segments_dict, slide_width, slide_height, company_name)
